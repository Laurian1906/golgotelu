"""
* File name: ppt.py
* Description:
This script will execute the extraction of the lyrics from the PowerPoint
presentations songs. The AI will process those songs using the song_processor.py
"""

# IMPORTS
from pptx import Presentation
from pathlib import Path

class PptProcessor:
    def __init__(self, ppt_path: Path):
        self.__ppt_path = ppt_path

    def __process(self):
        text_runs = []
        prs = Presentation(str(self.__ppt_path))

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return text_runs

    @staticmethod
    def extract_lyrics(ppt_path: Path):
        """
        Extracts lyrics from the PowerPoint file.
        :param ppt_path: Path to the PowerPoint presentation file.
        :return: A string containing the extracted lyrics from all slides.
        """

        # Get the path to the PowerPoint presentation
        processor = PptProcessor(ppt_path)

        # Return the extracted text
        return processor.__process()